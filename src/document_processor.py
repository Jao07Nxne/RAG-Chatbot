"""
Document Processor for Thai RAG Chatbot
รองรับการอ่านและประมวลผลเอกสารภาษาไทย
รองรับ Dynamic Chunking Strategy
"""

import os
import re
from typing import List, Dict, Any
from pathlib import Path
import chardet

# Document reading libraries
import PyPDF2
from docx import Document
from pptx import Presentation

# Thai language processing
import pythainlp
from pythainlp.tokenize import word_tokenize
from pythainlp.corpus.common import thai_words

# LangChain text processing
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document as LangChainDocument

# Dynamic Chunking
from dynamic_text_splitter import DynamicTextSplitter
from content_classifier import ContentType


class ThaiDocumentProcessor:
    """คลาสสำหรับประมวลผลเอกสารภาษาไทย พร้อม Dynamic Chunking"""
    
    def __init__(self, chunk_size: int = 1000, chunk_overlap: int = 200, use_dynamic_chunking: bool = True):
        """
        Args:
            chunk_size: ขนาดของแต่ละ chunk (characters) - ใช้เมื่อ use_dynamic_chunking=False
            chunk_overlap: ส่วนที่ซ้อนทับระหว่าง chunks - ใช้เมื่อ use_dynamic_chunking=False
            use_dynamic_chunking: ใช้ Dynamic Chunking Strategy หรือไม่ (แนะนำ True)
        """
        self.chunk_size = chunk_size
        self.chunk_overlap = chunk_overlap
        self.use_dynamic_chunking = use_dynamic_chunking
        
        if use_dynamic_chunking:
            # ใช้ Dynamic Text Splitter (แนะนำ!)
            print("✨ เปิดใช้งาน Dynamic Chunking Strategy")
            self.dynamic_splitter = DynamicTextSplitter()
            self.text_splitter = None  # ไม่ใช้ fixed splitter
        else:
            # ใช้ Fixed Text Splitter (แบบเดิม)
            print(f"⚠️ ใช้ Fixed Chunking ({chunk_size}/{chunk_overlap})")
            self.dynamic_splitter = None
            self.text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=chunk_size,
                chunk_overlap=chunk_overlap,
                length_function=len,
                separators=[
                    "\n\n3.1.4",   # 🔥 จับ section "แผนการศึกษา" ให้แยก chunk
                    "\n\n3.1.5",   # จับ section ถัดไป
                    "\n\nปีที่ ",   # จับหัวข้อปี (เช่น "ปีที่ 1", "ปีที่ 2")
                    "\n\nภาคการศึกษาที่ ",  # จับหัวข้อภาค
                    "\n\n\n",      # Multiple paragraph breaks
                    "\n\n",        # Paragraph breaks
                    "\n",          # Line breaks
                    "。",          # Full stop (Thai/Asian)
                    ".",           # Sentence endings
                    "!",           # Exclamation
                    "?",           # Question
                    ";",           # Semicolon
                    ",",           # Comma
                    " ",           # Spaces
                    ""             # Characters (ใช้เมื่อจำเป็น)
                ]
            )
    
    def detect_encoding(self, file_path: str) -> str:
        """ตัวอย่างการตรวจสอบ encoding ของไฟล์"""
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                return result['encoding'] or 'utf-8'
        except Exception:
            return 'utf-8'
    
    def read_text_file(self, file_path: str) -> str:
        """อ่านไฟล์ text ธรรมดา"""
        encoding = self.detect_encoding(file_path)
        try:
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except UnicodeDecodeError:
            # ลองใช้ encoding อื่น
            for enc in ['utf-8', 'cp874', 'iso-8859-11']:
                try:
                    with open(file_path, 'r', encoding=enc) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue
            raise ValueError(f"ไม่สามารถอ่านไฟล์ {file_path} ได้")
    
    def read_pdf_file(self, file_path: str) -> str:
        """อ่านไฟล์ PDF ด้วยการปรับปรุงคุณภาพข้อความ"""
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    
                    # ลบช่องว่างพิเศษระหว่างตัวอักษรไทย (ปัญหาหลักของ PDF)
                    # "ร ้ า น ง า น" → "ร้านงาน"
                    page_text = re.sub(r'([\u0E00-\u0E7F])\s+([\u0E00-\u0E7F])', r'\1\2', page_text)
                    
                    # ลบช่องว่างระหว่างตัวอักษรไทยกับสระ/วรรณยุกต์
                    # "น ้ อ ย" → "น้อย"
                    page_text = re.sub(r'([\u0E00-\u0E7F])\s+([\u0E31-\u0E3A\u0E47-\u0E4E])', r'\1\2', page_text)
                    
                    # ลบช่องว่างก่อนสระบน/ล่าง
                    page_text = re.sub(r'\s+([\u0E31-\u0E3A\u0E47-\u0E4E])', r'\1', page_text)
                    
                    text += page_text + "\n"
                    
                    # Debug: แสดงตัวอย่างข้อความที่อ่านได้
                    if page_num == 0:
                        print(f"📄 ตัวอย่างข้อความหน้า 1: {page_text[:200]}")
            
            return text
        except Exception as e:
            raise ValueError(f"ไม่สามารถอ่านไฟล์ PDF {file_path}: {str(e)}")
    
    def read_docx_file(self, file_path: str) -> str:
        """อ่านไฟล์ Word (.docx)"""
        try:
            doc = Document(file_path)
            text = ""
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"ไม่สามารถอ่านไฟล์ Word {file_path}: {str(e)}")
    
    def read_pptx_file(self, file_path: str) -> str:
        """อ่านไฟล์ PowerPoint (.pptx)"""
        try:
            prs = Presentation(file_path)
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        except Exception as e:
            raise ValueError(f"ไม่สามารถอ่านไฟล์ PowerPoint {file_path}: {str(e)}")
    
    def clean_thai_text(self, text: str) -> str:
        """ทำความสะอาดข้อความภาษาไทย (ปรับปรุงสำหรับ PDF ที่มีปัญหา)"""
        # ขั้นตอนที่ 1: ลบช่องว่างพิเศษระหว่างตัวอักษรไทย
        # "ก ว า ย" → "กวาย"
        text = re.sub(r'([\u0E00-\u0E7F])\s+([\u0E00-\u0E7F])', r'\1\2', text)
        
        # ขั้นตอนที่ 2: ลบช่องว่างระหว่างตัวอักษรกับสระ/วรรณยุกต์
        # "น ้ อ ย" → "น้อย", "ร ั บ" → "รับ"
        text = re.sub(r'([\u0E00-\u0E7F])\s+([\u0E31-\u0E3A\u0E47-\u0E4E])', r'\1\2', text)
        text = re.sub(r'\s+([\u0E31-\u0E3A\u0E47-\u0E4E])\s+', r'\1', text)
        
        # ขั้นตอนที่ 3: 🔥 แก้ช่องว่างในรหัสวิชา ก่อนแก้อย่างอื่น
        # "0550 6231" → "05506231", "0550  6231" → "05506231"
        text = re.sub(r'\b(\d{4})\s+(\d{4})\b', r'\1\2', text)
        text = re.sub(r'\b(\d{2})\s+(\d{2})\s+(\d{4})\b', r'\1\2\3', text)
        
        # ขั้นตอนที่ 4: แก้ตัวเลข/อักขระพิเศษที่ติดคำไทย
        text = re.sub(r'คอมพิวเตอร4', 'คอมพิวเตอร์', text)
        text = re.sub(r'([\u0E00-\u0E7F]+)4', r'\1์', text)  # แปลง "คำ4" → "คำ์"
        text = re.sub(r'([\u0E00-\u0E7F]+)e', r'\1ี', text)  # แปลง "ปe" → "ปี"
        text = re.sub(r'([\u0E00-\u0E7F]+)F', r'\1้', text)  # แปลง "หนF" → "หน้"
        text = re.sub(r'([\u0E00-\u0E7F]+)Q', r'\1้', text)  # แปลง "กQ" → "ก้"
        
        # เพิ่มตัวอักษรพิเศษอื่นๆ ที่พบในเอกสาร
        text = re.sub(r'เปbด', 'เปิด', text)
        text = re.sub(r'ไม้น้อยกว้า', 'ไม่น้อยกว่า', text)
        text = re.sub(r'ได้แก้', 'ได้แก่', text)
        text = re.sub(r'ฝå', 'ฝึ', text)
        text = re.sub(r'พระจอมเกล<าเจ<า', 'พระจอมเกล้าเจ้า', text)
        
        # ขั้นตอนที่ 5: แก้คำที่พบบ่อยซึ่งมักมีปัญหาใน PDF
        common_fixes = {
            r'ห น ่ ว ย': 'หน่วย',
            r'ห น ่ วย': 'หน่วย',
            r'หนวยกิต': 'หน่วยกิต',
            r'ก ิ ต': 'กิต',
            r'ก ิต': 'กิต',
            r'ห ล ั ก ส ู ต ร': 'หลักสูตร',
            r'ห ล ั กส ู ตร': 'หลักสูตร',
            r'ว ิ ท ย า': 'วิทยา',
            r'ว ิทย า': 'วิทยา',
            r'ศ ึ ก ษ า': 'ศึกษา',
            r'ศ ึกษ า': 'ศึกษา',
            r'ภ า ค': 'ภาค',
            r'เท อ ม': 'เทอม',
            r'ป ร ะ จ ำ': 'ประจำ',
            r'เกลQาเจQา': 'เกล้าเจ้า',
            r'วFา': 'ว่า',
            r'ดQวย': 'ด้วย',
            # 🔥 เพิ่มคำที่พบจากการทดสอบ
            r'คดีศาสตร์ดิลิลธิล': 'คณิตศาสตร์ดิสครีต',
            r'คดีศาสตร์': 'คณิตศาสตร์',
            r'ดิลิลธิล': 'ดิสครีต',
            r'ปฏิสัมพันระ': 'ปฏิสัมพันธ์',
            r'สวัดกรรม': 'วิศวกรรม',
            r'สวัด': 'วิศว',
        }
        
        for pattern, replacement in common_fixes.items():
            text = re.sub(pattern, replacement, text)
        
        # ขั้นตอนที่ 6: ลบ whitespace ที่ไม่จำเป็น (มากกว่า 1 ช่องว่าง)
        text = re.sub(r' {2,}', ' ', text)
        
        # ขั้นตอนที่ 7: ลบอักขระพิเศษที่ไม่ต้องการ (แต่เก็บอักขระไทยและตัวเลข)
        text = re.sub(r'[^\u0E00-\u0E7F\w\s\.\,\!\?\;\:\-\(\)\"\'\/\n]', '', text)
        
        # ขั้นตอนที่ 8: ลบบรรทัดว่างซ้ำ (มากกว่า 2 บรรทัด)
        text = re.sub(r'\n{3,}', '\n\n', text)
        
        # ขั้นตอนที่ 9: ลบช่องว่างต้นและท้ายบรรทัด
        text = '\n'.join(line.strip() for line in text.split('\n'))
        
        return text.strip()
    
    def preprocess_text(self, text: str) -> str:
        """ประมวลผลข้อความก่อนการแบ่ง chunks"""
        # ทำความสะอาด
        text = self.clean_thai_text(text)
        
        # ตัด word ภาษาไทยเพื่อปรับปรุงการแบ่ง chunk
        # (อาจจะใช้หรือไม่ใช้ก็ได้ ขึ้นอยู่กับความต้องการ)
        
        return text
    
    def read_document(self, file_path: str) -> str:
        """อ่านเอกสารตามประเภทไฟล์"""
        file_path = Path(file_path)
        file_extension = file_path.suffix.lower()
        
        if file_extension == '.txt':
            return self.read_text_file(str(file_path))
        elif file_extension == '.pdf':
            return self.read_pdf_file(str(file_path))
        elif file_extension in ['.docx', '.doc']:
            return self.read_docx_file(str(file_path))
        elif file_extension in ['.pptx', '.ppt']:
            return self.read_pptx_file(str(file_path))
        else:
            # ลองอ่านเป็น text file
            try:
                return self.read_text_file(str(file_path))
            except:
                raise ValueError(f"ไม่รองรับไฟล์ประเภท {file_extension}")
    
    def process_document(self, file_path: str, metadata: Dict[str, Any] = None) -> List[LangChainDocument]:
        """
        ประมวลผลเอกสารและแบ่งเป็น chunks (รองรับ Dynamic Chunking)
        
        Args:
            file_path: เส้นทางไฟล์
            metadata: ข้อมูลเพิ่มเติมของเอกสาร
            
        Returns:
            List[LangChainDocument]: รายการ chunks พร้อม metadata
        """
        # อ่านเอกสาร
        text = self.read_document(file_path)
        
        # ประมวลผลข้อความ
        processed_text = self.preprocess_text(text)
        
        if not processed_text.strip():
            print("⚠️ ไม่พบข้อความในเอกสาร")
            return []
        
        # แบ่ง chunks ตาม strategy
        if self.use_dynamic_chunking:
            # ใช้ Dynamic Chunking
            print("\n🎯 เริ่มแบ่ง chunks ด้วย Dynamic Strategy...")
            chunks, content_type = self.dynamic_splitter.split_text(processed_text)
        else:
            # ใช้ Fixed Chunking (แบบเดิม)
            print(f"\n📝 เริ่มแบ่ง chunks (Fixed: {self.chunk_size}/{self.chunk_overlap})...")
            chunks = self.text_splitter.split_text(processed_text)
            content_type = "general"  # Default
        
        # สร้าง LangChain Documents พร้อม metadata ที่ช่วยในการค้นหา
        documents = []
        file_name = Path(file_path).name
        
        base_metadata = {
            "source": file_path,
            "filename": file_name,
            "total_chunks": len(chunks)
        }
        
        if metadata:
            base_metadata.update(metadata)
        
        # เพิ่ม content_type ถ้าใช้ Dynamic Chunking
        if self.use_dynamic_chunking:
            base_metadata["content_type"] = content_type
        
        for i, chunk in enumerate(chunks):
            chunk_metadata = base_metadata.copy()
            chunk_metadata.update({
                "chunk_index": i,
                "chunk_size": len(chunk)
            })
            
            # 🔍 เพิ่ม metadata เฉพาะสำหรับการค้นหาที่แม่นยำขึ้น
            # ตรวจจับรหัสวิชา (8 หลัก เช่น 05506232, 90641001)
            course_codes = re.findall(r'\b\d{8}\b', chunk)
            if course_codes:
                chunk_metadata["course_codes"] = ",".join(list(set(course_codes)))
            
            # 🔥 ตรวจจับชั้นปีและภาคการศึกษา (ปรับ regex ให้จับได้หลายรูปแบบ)
            # รูปแบบ 1: "ปีที่ 1" หรือ "ปีที่1"
            year_patterns = [
                r'ปีที่\s*(\d+)',
                r'ชั้นปี\s*(\d+)',
                r'แผนการศึกษา\s*ปีที่\s*(\d+)',
            ]
            year_found = None
            for pattern in year_patterns:
                match = re.search(pattern, chunk)
                if match:
                    year_found = match.group(1)
                    break
            
            # รูปแบบ 2: "ภาคการศึกษาที่ 1" หรือ "ภาค 1"
            semester_patterns = [
                r'ภาคการศึกษาที่\s*(\d+)',
                r'ภาค\s*(\d+)',
                r'เทอม\s*(\d+)',
            ]
            semester_found = None
            for pattern in semester_patterns:
                match = re.search(pattern, chunk)
                if match:
                    semester_found = match.group(1)
                    break
            
            if year_found:
                chunk_metadata["year"] = year_found
                # ถ้ามีทั้งปีและภาค → tag ว่าเป็น curriculum table
                if semester_found:
                    chunk_metadata["semester"] = semester_found
                    chunk_metadata["is_curriculum_table"] = "yes"
                    
                    # นับจำนวนวิชาในตาราง
                    if course_codes:
                        chunk_metadata["course_count"] = len(course_codes)
                    
                    # ดึงจำนวนหน่วยกิตรวม
                    total_match = re.search(r'รวม\s+(\d+)\s+หน่วยกิต', chunk)
                    if total_match:
                        chunk_metadata["total_credits"] = total_match.group(1)
                    
                    print(f"   🎯 พบตารางหลักสูตร: ปี {year_found} ภาค {semester_found} "
                          f"({len(course_codes)} วิชา) ใน chunk {i}")
            elif semester_found:
                chunk_metadata["semester"] = semester_found
            
            # ตรวจจับหน่วยกิต
            credit_matches = re.findall(r'(\d+)\s*หน่วยกิต', chunk)
            if credit_matches:
                chunk_metadata["has_credits"] = "yes"
            
            # เพิ่มคำสำคัญจาก chunk เป็น preview
            chunk_metadata["preview"] = chunk[:150].strip().replace('\n', ' ')
            
            documents.append(LangChainDocument(
                page_content=chunk,
                metadata=chunk_metadata
            ))
        
        print(f"✅ แบ่ง chunks เสร็จสิ้น: {len(documents)} chunks")
        
        return documents
    
    def process_multiple_documents(self, file_paths: List[str]) -> List[LangChainDocument]:
        """ประมวลผลหลายเอกสารพร้อมกัน"""
        all_documents = []
        
        for file_path in file_paths:
            try:
                documents = self.process_document(file_path)
                all_documents.extend(documents)
                print(f"✅ ประมวลผลไฟล์ {Path(file_path).name} เสร็จสิ้น - {len(documents)} chunks")
            except Exception as e:
                print(f"❌ ไม่สามารถประมวลผลไฟล์ {Path(file_path).name}: {str(e)}")
        
        return all_documents
    
    def get_text_stats(self, text: str) -> Dict[str, Any]:
        """สถิติของข้อความ"""
        words = word_tokenize(text, engine='newmm')
        thai_words_count = len([w for w in words if re.match(r'[\u0E00-\u0E7F]+', w)])
        
        return {
            "total_characters": len(text),
            "total_words": len(words),
            "thai_words": thai_words_count,
            "lines": len(text.split('\n')),
            "estimated_chunks": (len(text) // self.chunk_size) + 1
        }


# ตัวอย่างการใช้งาน
if __name__ == "__main__":
    processor = ThaiDocumentProcessor(chunk_size=1000, chunk_overlap=200)
    
    # ทดสอบอ่านไฟล์
    try:
        # สร้างไฟล์ทดสอบ
        test_text = """
        ตัวอย่างเอกสารภาษาไทยสำหรับทดสอบ RAG Chatbot
        
        ประเทศไทยมีชื่อเป็นทางการว่า ราชอาณาจักรไทย เป็นประเทศในเอเชียตะวันออกเฉียงใต้ 
        มีพรมแดนติดกับพม่า ลาว กัมพูชา และมาเลเซีย มีเนื้อที่ประมาณ 513,120 ตารางกิโลเมตร
        
        กรุงเทพมหานครเป็นเมืองหลวงและเมืองที่ใหญ่ที่สุดของประเทศไทย 
        ประชากรประมาณ 70 ล้านคน ใช้ภาษาไทยเป็นภาษาราชการ
        
        ระบบการเมืองเป็นราชาธิปไตยภายใต้รัธรรมนูญ โดยมีพระมหากษัตริย์เป็นประมุข
        """
        
        with open("test_thai.txt", "w", encoding="utf-8") as f:
            f.write(test_text)
        
        # ทดสอบการประมวลผล
        docs = processor.process_document("test_thai.txt")
        
        print(f"จำนวน chunks: {len(docs)}")
        for i, doc in enumerate(docs):
            print(f"\nChunk {i+1}:")
            print(f"Content: {doc.page_content[:100]}...")
            print(f"Metadata: {doc.metadata}")
        
        # ลบไฟล์ทดสอบ
        os.remove("test_thai.txt")
        
    except Exception as e:
        print(f"เกิดข้อผิดพลาด: {e}")